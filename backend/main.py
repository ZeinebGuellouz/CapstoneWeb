
from fastapi import FastAPI, File, UploadFile, HTTPException, Depends, Request
from fastapi.middleware.cors import CORSMiddleware
from fastapi.staticfiles import StaticFiles
from fastapi.responses import FileResponse
from pathlib import Path
from typing import List
import shutil
import os
import pythoncom
import comtypes.client
from pptx import Presentation
from pdf2image import convert_from_path
import pdfplumber
from pydantic import BaseModel
import requests
import firebase_admin
from firebase_admin import credentials, auth, firestore
from fastapi import Header , Query



# ========== FASTAPI SETUP ==========
app = FastAPI()

# Directory to store slide images
SLIDE_DIR = Path("slides")
SLIDE_DIR.mkdir(exist_ok=True)

# Directory to store uploaded files
UPLOAD_DIR = Path("uploads")
UPLOAD_DIR.mkdir(exist_ok=True)

os.environ["GOOGLE_APPLICATION_CREDENTIALS"] = "presentpro-b7e76-firebase-adminsdk-fbsvc-f540aa32d5.json"

cred = credentials.Certificate("presentpro-b7e76-firebase-adminsdk-fbsvc-f540aa32d5.json")
firebase_admin.initialize_app(cred)
db = firestore.Client()

# ✅ Function to verify Firebase Token
def verify_firebase_token(token: str):
    try:
        decoded_token = auth.verify_id_token(token)
        return decoded_token  # Returns user data
    except Exception as e:
        print(f"Token verification failed: {e}")
        return None

# Enable CORS
app.add_middleware(
    CORSMiddleware,
    allow_origins=["*"],  
    allow_credentials=True,
    allow_methods=["*"],
    allow_headers=["*"],
)

# ========== SIGNUP (Handled by Firebase) ==========
class UserSignup(BaseModel):
    email: str
    password: str

@app.post("/signup")
def signup(user: UserSignup):
    try:
        # ✅ Create user in Firebase
        firebase_user = auth.create_user(email=user.email, password=user.password)
        return {"message": "User registered successfully", "uid": firebase_user.uid}
    except Exception as e:
        raise HTTPException(status_code=400, detail=str(e))

# ✅ API Route for OAuth Login (Google & Facebook)
@app.post("/oauth-login")
def oauth_login(authorization: str = Header(...)):
    if not authorization or not authorization.startswith("Bearer "):
        raise HTTPException(status_code=401, detail="Missing or invalid token format")

    token = authorization.split("Bearer ")[1]  # ✅ Extract the token
    user_data = verify_firebase_token(token)

    if not user_data:
        raise HTTPException(status_code=401, detail="Invalid Firebase Token")

    return {"message": "OAuth login successful", "email": user_data.get("email")}

# ✅ PASSWORD RESET (Handled by Firebase)
@app.post("/reset-password")
def reset_password(email: str):
    try:
        auth.generate_password_reset_link(email)
        return {"message": "Password reset email sent successfully"}
    except firebase_admin.exceptions.FirebaseError as e:
        raise HTTPException(status_code=400, detail=str(e))


# ========== DATABASE DEPENDENCIES ==========

def clear_slide_directory(user_id: str, presentation_id: str) -> Path:
    slide_dir = SLIDE_DIR / user_id / presentation_id
    if slide_dir.exists():
        shutil.rmtree(slide_dir)
    slide_dir.mkdir(parents=True, exist_ok=True)
    return slide_dir


def process_pdf(file_path: Path, output_dir: Path) -> List[str]:
    try:
        slides = convert_from_path(file_path, dpi=150, output_folder=output_dir, fmt="png")
        return [f"/uploads/slides/{output_dir.relative_to(SLIDE_DIR)}/{Path(s.filename).name}" for s in slides]

    except Exception as e:
        print(f"Error processing PDF: {e}")
        return []

def extract_text_from_pptx(file_path: Path) -> List[str]:
    try:
        prs = Presentation(file_path)
        slide_texts = []
        for slide in prs.slides:
            text = []
            for shape in slide.shapes:
                if shape.has_text_frame:
                    text.append(shape.text)
            slide_texts.append("\n".join(text))
        return slide_texts
    except Exception as e:
        print(f"Error extracting text from PPTX: {e}")
        return []

def extract_text_from_pdf(file_path: Path) -> List[str]:
    try:
        slide_texts = []
        with pdfplumber.open(file_path) as pdf:
            for page in pdf.pages:
                slide_texts.append(page.extract_text() or "")
        return slide_texts
    except Exception as e:
        print(f"Error extracting text from PDF: {e}")
        return []

def process_pptx(file_path: Path, output_dir: Path) -> List[str]:
    pythoncom.CoInitialize()
    pptx_file = os.path.abspath(file_path)
    output_folder = os.path.abspath(output_dir)
    os.makedirs(output_folder, exist_ok=True)

    try:
        powerpoint = comtypes.client.CreateObject("PowerPoint.Application")
        presentation = powerpoint.Presentations.Open(pptx_file, WithWindow=False)

        slide_paths = []
        for i, slide in enumerate(presentation.Slides):
            output_image = os.path.join(output_folder, f"slide_{i + 1}.png")
            slide.Export(output_image, "PNG", 1280, 720)
            slide_paths.append(f"/uploads/slides/{output_dir.relative_to(SLIDE_DIR)}/slide_{i + 1}.png")

        presentation.Close()
        powerpoint.Quit()
        return slide_paths
    except Exception as e:
        print(f"Error processing PPTX: {e}")
        return []



# ========== Pydantic Models ==========
class SpeechRequest(BaseModel):
    presentation_id: str
    slide_number: int
    generated_speech: str
    voice_tone: str = "Formal"
    speed: float = 1.0
    pitch: float = 1.0

# ========== API to Generate Speech ==========
@app.post("/generate_speech/")
async def generate_speech(request: Request):
    try:
        body = await request.json()
        previous_slides = body.get("previous_slides", [])
        current_slide = body.get("current_slide", {})
        slide_index = body.get("slide_index", 0)
        voice_tone = body.get("voice_tone", "Formal")  # Default to "Formal" if not provided

        if not current_slide:
            return {"error": "No current slide provided"}

        first_slide_text = previous_slides[0]["text"] if previous_slides else current_slide["text"]
        language = "French" if any(word in first_slide_text.lower() for word in ["le", "la", "les", "un", "une", "est", "avec"]) else "English"

        prompt = (
             f"You are an AI assistant generating a **{voice_tone.lower()}**, natural-sounding speech for a {language} presentation. "
             f"Your speech should be short, engaging, and directly relevant to the slide, keeping it in {language}.\n"
             f"Avoid unnecessary introductions like 'Ladies and gentlemen' or 'Mesdames et messieurs' and repetitive information.\n\n"
)


        if previous_slides:
            prompt += f"The previous slides covered (in {language}):\n"
            for i, slide in enumerate(previous_slides):
                prompt += f"- Slide {i+1}: {slide['text']}\n"
            prompt += "\n"

        prompt += (
            f"Now, generate **a short, clear, and engaging speech** ONLY for **Slide {slide_index+1}** in **{language}**:\n"
            f"{current_slide['text']}\n\n"
            "Keep the response concise (**3-5 sentences**), avoiding repetition of prior slides."
        )

        response = requests.post(
            "http://localhost:11434/api/generate",
            json={
                "model": "mistral",
                "prompt": prompt,
                "stream": False,
                "max_tokens": 150
            }
        )

        if response.status_code != 200:
            return {"error": f"Ollama API Error: {response.text}"}

        data = response.json()
        generated_speech = data.get("response", "").strip()

        if not generated_speech:
            return {"error": "Empty response from Mistral"}

        return {"speech": generated_speech}

    except requests.exceptions.ConnectionError:
        return {"error": "Failed to connect to Ollama. Ensure it's running."}

    except Exception as e:
        return {"error": f"Unexpected error: {str(e)}"}

@app.get("/uploads/slides/{path:path}")
async def get_slide(path: str):
    file_path = SLIDE_DIR / path
    if not file_path.exists():
        raise HTTPException(status_code=404, detail="File not found.")
    return FileResponse(file_path, headers={"Cache-Control": "no-cache, no-store, must-revalidate"})

@app.post("/save_speech")
async def save_speech(request: Request, authorization: str = Header(...)):
    if not authorization.startswith("Bearer "):
        raise HTTPException(status_code=401, detail="Missing or invalid authorization header")

    id_token = authorization.split("Bearer ")[1]
    try:
        decoded_token = auth.verify_id_token(id_token)
        user_id = decoded_token["uid"]
    except:
        raise HTTPException(status_code=401, detail="Invalid Firebase ID token")

    body = await request.json()

    presentation_id = body.get("presentation_id")
    slide_number = str(body.get("slide_number"))

    if not presentation_id or not slide_number:
        raise HTTPException(status_code=400, detail="Missing presentation ID or slide number")

    slide_ref = (
        db.collection("users")
        .document(user_id)
        .collection("presentations")
        .document(presentation_id)
        .collection("slides")
        .document(slide_number)
    )

    # ✅ Update speech, pitch, tone, text (optional), and lastModifiedAt
    update_data = {
        "speech": body.get("generated_speech"),
        "pitch": body.get("pitch", 1.0),
        "speed": body.get("speed", 1.0),
        "voice_tone": body.get("voice_tone", "Formal"),
        "text": body.get("text"),  # optional if user edited it
        "lastModifiedAt": firestore.SERVER_TIMESTAMP
    }

    # Remove any None values
    update_data = {k: v for k, v in update_data.items() if v is not None}

    slide_ref.update(update_data)
    return {"message": "Speech data saved successfully to Firestore"}



@app.post("/upload/")
async def upload_file(
    file: UploadFile = File(...),
    authorization: str = Header(...)
):
    if not authorization or not authorization.startswith("Bearer "):
        raise HTTPException(status_code=401, detail="Missing or invalid authorization header")
    
    id_token = authorization.split("Bearer ")[1]
    try:
        decoded_token = auth.verify_id_token(id_token)
        user_id = decoded_token["uid"]
    except:
        raise HTTPException(status_code=401, detail="Invalid Firebase ID token")

    file_ext = file.filename.split(".")[-1].lower()
    if file_ext not in ["pptx", "pdf"]:
        raise HTTPException(status_code=400, detail="Unsupported file format")

    # ✅ Save the file to disk
    file_path = UPLOAD_DIR / file.filename
    with file_path.open("wb") as buffer:
        shutil.copyfileobj(file.file, buffer)

    # ✅ Generate unique ID after file is saved (to get correct mtime)
    presentation_id = f"{int(file_path.stat().st_mtime * 1000)}_{file.filename.replace(' ', '_')}"

    # ✅ Clear/create a fresh directory for this presentation
    slide_dir = clear_slide_directory(user_id, presentation_id)

    try:
        if file_ext == "pptx":
            slide_images = process_pptx(file_path, slide_dir)
            slide_texts = extract_text_from_pptx(file_path)
        else:
            slide_images = process_pdf(file_path, slide_dir)
            slide_texts = extract_text_from_pdf(file_path)
    except Exception as e:
        raise HTTPException(status_code=500, detail=f"Processing failed: {e}")

    if not slide_images:
        raise HTTPException(status_code=500, detail="No slides extracted")

    # ✅ Save metadata to Firestore (thumbnail path is user/pres specific now)
    presentation_doc = db.collection("users").document(user_id).collection("presentations").document(presentation_id)
    thumbnail_relative_path = f"/uploads/slides/{user_id}/{presentation_id}/{Path(slide_images[0]).name}"
    presentation_doc.set({
        "fileName": file.filename,
        "createdAt": firestore.SERVER_TIMESTAMP,
        "thumbnailUrl": f"http://127.0.0.1:8000{thumbnail_relative_path}",
    })

    slides_collection = presentation_doc.collection("slides")
    for idx, (img_path, text) in enumerate(zip(slide_images, slide_texts), start=1):
        slides_collection.document(str(idx)).set({
            "slideNumber": idx,
            "text": text,
            "image": f"http://127.0.0.1:8000{img_path}"
        })

    return {
        "presentationId": presentation_id,
        "filename": file.filename,
        "slides": [
            {"image": f"http://127.0.0.1:8000{path}", "text": text}
            for path, text in zip(slide_images, slide_texts)
        ]
    }

@app.get("/my_presentations")
async def get_user_presentations(authorization: str = Header(...)):
    if not authorization.startswith("Bearer "):
        raise HTTPException(status_code=401, detail="Invalid auth header")

    id_token = authorization.split("Bearer ")[1]
    try:
        decoded_token = auth.verify_id_token(id_token)
        user_id = decoded_token["uid"]
    except:
        raise HTTPException(status_code=401, detail="Invalid Firebase ID token")

    presentations_ref = db.collection("users").document(user_id).collection("presentations")
    presentations = presentations_ref.stream()

    result = []
    for doc in presentations:
        data = doc.to_dict()
        presentation_id = doc.id

        # Convert Firestore timestamp to frontend-safe format
        created_at = data.get("createdAt")
        timestamp = {
            "_seconds": int(created_at.timestamp())
        } if created_at else None

        result.append({
            "id": presentation_id,
            "fileName": data.get("fileName"),
            "createdAt": timestamp,
            "thumbnailUrl": data.get("thumbnailUrl")
        })

    return result

@app.delete("/delete_presentation")
async def delete_presentation(
    presentation_id: str = Query(...),
    authorization: str = Header(...)
):
    if not authorization.startswith("Bearer "):
        raise HTTPException(status_code=401, detail="Invalid auth header")

    id_token = authorization.split("Bearer ")[1]
    try:
        decoded_token = auth.verify_id_token(id_token)
        user_id = decoded_token["uid"]
    except:
        raise HTTPException(status_code=401, detail="Invalid Firebase ID token")

    try:
        presentation_ref = (
            db.collection("users")
            .document(user_id)
            .collection("presentations")
            .document(presentation_id)
        )

        # Delete all slides in the subcollection
        slides = presentation_ref.collection("slides").stream()
        for slide in slides:
            slide.reference.delete()

        # Delete the presentation document
        presentation_ref.delete()

        return {"message": "Presentation deleted successfully"}
    except Exception as e:
        raise HTTPException(status_code=500, detail=str(e))
class BatchDeleteRequest(BaseModel):
    presentation_ids: List[str]

@app.post("/delete_presentations_batch")
async def delete_presentations_batch(
    request: BatchDeleteRequest,
    authorization: str = Header(...)
):
    if not authorization.startswith("Bearer "):
        raise HTTPException(status_code=401, detail="Invalid auth header")

    id_token = authorization.split("Bearer ")[1]
    try:
        decoded_token = auth.verify_id_token(id_token)
        user_id = decoded_token["uid"]
    except:
        raise HTTPException(status_code=401, detail="Invalid Firebase ID token")

    try:
        for pid in request.presentation_ids:
            pres_ref = db.collection("users").document(user_id).collection("presentations").document(pid)

            # Delete all slides
            slides = pres_ref.collection("slides").stream()
            for slide in slides:
                slide.reference.delete()

            # Delete the presentation itself
            pres_ref.delete()

        return {"message": "Batch deletion successful"}
    except Exception as e:
        raise HTTPException(status_code=500, detail=str(e))


@app.get("/")
def root():
    return {"message": "Backend is running!"}